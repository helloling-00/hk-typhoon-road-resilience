"""
Generate all 13 defense PPT slides in one batch.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

THESIS = "/Users/helloling/workspace/thesis"

C_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
C_MID    = RGBColor(0x16, 0x21, 0x3e)
C_ACCENT = RGBColor(0x0f, 0x8b, 0x8d)
C_LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
C_YELLOW = RGBColor(0xff, 0xc8, 0x00)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_GRAY   = RGBColor(0xaa, 0xaa, 0xaa)
C_RED    = RGBColor(0xe0, 0x52, 0x52)

# ── Helpers ──────────────────────────────────────────────────────────────────
def to_rgb(c):
    """Accept hex string or RGBColor, return RGBColor."""
    if isinstance(c, RGBColor): return c
    if isinstance(c, str) and c.startswith("#"):
        c = c.lstrip("#")
        return RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
    return c

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill: shape.fill.solid(); shape.fill.fore_color.rgb = to_rgb(fill)
    else: shape.fill.background()
    if line: shape.line.color.rgb = to_rgb(line); shape.line.width = Pt(1)
    else: shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, color=C_WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    return txb

def img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def top_bar(slide, num, tag, title, tag_color=C_ACCENT):
    rect(slide, 0, 0, 13.33, 0.75, fill=C_MID)
    txt(slide, f"{num:02d}", 0.15, 0.12, 0.5, 0.5, size=11, color=C_GRAY)
    tag_w = max(1.5, len(tag)*0.22)
    rect(slide, 0.5, 0.17, tag_w, 0.40, fill=tag_color)
    txt(slide, tag, 0.55, 0.18, tag_w-0.1, 0.38,
        size=10, color=C_WHITE if tag_color != C_YELLOW else C_DARK,
        bold=True, align=PP_ALIGN.CENTER)
    txt(slide, title, 0.5+tag_w+0.15, 0.10, 12.5-tag_w, 0.58, size=22, color=C_WHITE, bold=True)

def bot_line(slide, text):
    rect(slide, 0, 7.28, 13.33, 0.04, fill=C_ACCENT)
    txt(slide, text, 0.3, 7.33, 12.0, 0.20, size=8, color=C_GRAY, italic=True)

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide1(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    # Decorative accent line
    rect(slide, 1.5, 2.8, 10.33, 0.06, fill=C_ACCENT)
    rect(slide, 1.5, 5.0, 10.33, 0.04, fill=C_YELLOW)

    txt(slide, "Pre-, During-, and Post-Event Resilience Analysis\nof Urban Road Networks under Typhoons",
        1.5, 1.8, 10.5, 1.2, size=30, color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
    txt(slide, "Evidence from TomTom Floating-Car Data in Hong Kong",
        1.5, 3.0, 10.5, 0.6, size=18, color=C_ACCENT, align=PP_ALIGN.LEFT)

    txt(slide, "MPhil Thesis Defense",
        1.5, 4.0, 10.5, 0.45, size=16, color=C_YELLOW, align=PP_ALIGN.LEFT)
    txt(slide, "Department of Urban Planning and Design  ·  The University of Hong Kong",
        1.5, 4.55, 10.5, 0.4, size=13, color=C_GRAY, align=PP_ALIGN.LEFT)
    txt(slide, "2026",
        1.5, 5.3, 10.5, 0.4, size=14, color=C_GRAY, align=PP_ALIGN.LEFT)

    rect(slide, 0, 7.28, 13.33, 0.04, fill=C_ACCENT)

    prs_ = Presentation(); prs_.slide_width=Inches(13.33); prs_.slide_height=Inches(7.5)
    return prs_

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: RESEARCH FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide2(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 2, "FRAMEWORK", "Research Questions", C_ACCENT)

    rqs = [
        ("RQ1", "Distribution",
         "What is the distribution of road speed deviations\nacross pre-, during-, and post-typhoon phases?",
         C_ACCENT),
        ("RQ2", "Spatiotemporal",
         "When and where do significant speed changes occur?\nWhat are the spatiotemporal patterns?",
         C_YELLOW),
        ("RQ3", "Determinants",
         "What grid-level factors explain why some areas\nexperience greater speed changes during typhoons?",
         RGBColor(0x4f,0xc3,0xf7)),
        ("RQ4", "Resilience",
         "How quickly do road networks recover after typhoon\npassage? What shapes recovery trajectories?",
         RGBColor(0x81,0xc7,0x84)),
    ]

    for i, (code, title, desc, color) in enumerate(rqs):
        y = 1.15 + i * 1.45
        rect(slide, 0.5, y, 12.3, 1.30, fill=C_MID, line=color)
        # Number badge
        rect(slide, 0.65, y+0.12, 1.2, 1.05, fill=color)
        txt(slide, code, 0.70, y+0.20, 1.1, 0.45, size=18, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, title, 0.70, y+0.60, 1.1, 0.35, size=8, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)
        # Description
        txt(slide, desc, 2.1, y+0.15, 8.5, 1.0, size=14, color=C_LIGHT)

    # Data note
    txt(slide, "Data: TomTom 30-min floating-car speed data  ·  3 typhoons (Mina, Yagiasha, Madum)  ·  95,167 physical roads  ·  500m grid analysis",
        0.5, 7.0, 12.0, 0.25, size=9, color=C_GRAY, italic=True)
    bot_line(slide, "Research framework  |  Four interconnected research questions")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: DATA & METHODS
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide3(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 3, "DATA", "Data & Methodology", C_ACCENT)

    # Left: Data
    rect(slide, 0.3, 0.95, 6.2, 3.0, fill=C_MID, line=C_ACCENT)
    txt(slide, "Data Sources", 0.5, 1.05, 5.8, 0.35, size=14, color=C_YELLOW, bold=True)

    data_items = [
        ("TomTom Floating-Car Data", "30-min relative speed, 2025 Sep–Oct, ~6,000–7,000 roads/slot"),
        ("Road Network", "95,167 physical roads (endpoint clustering), OSM road categories"),
        ("POI & Demographics", "OSM POI (10 categories), 2021 Census (500m buffer), estate-level population"),
        ("Typhoon Signals", "HKO signal timeline: S1→S3→S8→S9→S10 and downgrades"),
    ]
    y = 1.50
    for title, desc in data_items:
        rect(slide, 0.50, y, 0.07, 0.07, fill=C_ACCENT)
        txt(slide, title, 0.68, y-0.03, 5.5, 0.28, size=10, color=C_YELLOW, bold=True)
        txt(slide, desc, 0.68, y+0.22, 5.5, 0.35, size=8.5, color=C_LIGHT)
        y += 0.58

    # Right: Methods
    rect(slide, 6.85, 0.95, 6.2, 3.0, fill=C_MID, line=C_ACCENT)
    txt(slide, "Analytical Approach", 7.05, 1.05, 5.8, 0.35, size=14, color=C_YELLOW, bold=True)

    methods = [
        ("Y = Speed Deviation", "typhoon_speed − baseline_speed (by road type, day type, slot)"),
        ("Network Time Series", "Per-slot mean deviation across all matched roads"),
        ("500m Grid Aggregation", "Length-weighted road metrics → grid-level Y and X"),
        ("OLS Regression", "Grid-level: Y = f(POI, structure, demographics, signal)"),
    ]
    y = 1.50
    for title, desc in methods:
        rect(slide, 7.05, y, 0.07, 0.07, fill=C_YELLOW)
        txt(slide, title, 7.23, y-0.03, 5.5, 0.28, size=10, color=C_YELLOW, bold=True)
        txt(slide, desc, 7.23, y+0.22, 5.5, 0.35, size=8.5, color=C_LIGHT)
        y += 0.58

    # Bottom: Coverage figure
    rect(slide, 0.3, 4.15, 12.7, 2.85, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图19_路网覆盖分布图.png", 0.40, 4.20, 6.0, 2.75)
    img(slide, f"{THESIS}/图20_消失路演变图.png", 6.60, 4.20, 6.2, 2.75)

    bot_line(slide, "Data & Methods  |  95,167 roads  |  500m grid  |  OLS with length-weighted aggregation")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: OVERVIEW — Network-Level Time Series
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide4(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 4, "RQ1 · RQ2 · RQ4", "Overview: Network Response across All Typhoon Phases", C_ACCENT)

    rect(slide, 0.2, 0.85, 12.9, 4.9, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图42_台风时序网络偏差.png", 0.25, 0.90, 12.8, 4.80)

    # Takeaway strip
    rect(slide, 0.2, 5.85, 12.9, 1.08, fill=C_MID, line=C_ACCENT)
    takeaways = [
        ("Pre-event", "S8 is the behavioural trigger:\npre-S8 midday congestion dip"),
        ("During", "S8 triggers rapid network clearance\n(−0.010 → +0.045 in 5h)"),
        ("Post-event", "Recovery within ~100 min\nof All Clear; no rebound"),
    ]
    x = 0.45
    for phase, desc in takeaways:
        rect(slide, x, 5.95, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, phase, x+0.14, 5.93, 1.5, 0.30, size=11, color=C_YELLOW, bold=True)
        txt(slide, desc, x+0.14, 6.22, 3.8, 0.50, size=9.5, color=C_LIGHT)
        x += 4.3

    bot_line(slide, "Ragasa (Sep 2025)  |  30-min TomTom data  |  ~6,000 roads per slot  |  peak +0.046 at S10 morning")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: PRE-EVENT — S8 Is the True Behavioural Trigger
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide5(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 5, "RQ1 · RQ2", "Pre-Event: S8 Is the True Behavioural Trigger", C_YELLOW)

    # MAIN FIGURE: Mina vs Yagiasha comparison (图47) — spans full width
    rect(slide, 0.2, 0.85, 12.9, 3.15, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图47_速度形状对比_米娜vs叶加沙.png", 0.25, 0.90, 12.8, 3.05)
    txt(slide, "A  Mina (no S8): all-day elevated speeds, including midday  —  Yagiasha (S8+): midday dip ONLY on Sep 23 (pre-S8)",
        0.35, 0.87, 12.0, 0.25, size=8, color=C_ACCENT, bold=True)

    # BOTTOM ROW: Two panels — detailed pre-S8 + spatial
    # Left: Detailed pre-S8 surge
    rect(slide, 0.2, 4.15, 6.3, 3.0, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图43b_preS8_surge.png", 0.25, 4.20, 6.2, 2.90)
    txt(slide, "B  Yagiasha Sep 23: Pre-S8 dip → S8 clearance in detail", 0.35, 4.17, 6.0, 0.25,
        size=8, color=C_ACCENT, bold=True)

    # Right: Findings panel
    rect(slide, 6.85, 4.15, 6.25, 3.0, fill=C_MID, line=C_ACCENT)
    txt(slide, "Key Evidence", 7.05, 4.23, 5.8, 0.30, size=13, color=C_YELLOW, bold=True)

    evidence = [
        ("Mina (max S3, no S8)",
         "Sep 18–19 midday: speeds +0.005 to +0.015 above baseline\n→ No congestion surge; demand already suppressed"),
        ("Yagiasha (S8 raised 14:20)",
         "Sep 23 midday: speed dips to −0.010 at 13:00\n→ 31.3% roads slower, system-wide last-minute rush"),
        ("5-hour swing",
         "−0.010 (13:00) → +0.045 (19:30) after S8\n→ S8 is the behavioural tipping point, not S1 or S3"),
    ]
    y = 4.65
    for title, desc in evidence:
        rect(slide, 7.05, y+0.03, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, title, 7.25, y, 5.5, 0.25, size=10, color=C_ACCENT, bold=True)
        txt(slide, desc, 7.25, y+0.22, 5.5, 0.55, size=8.5, color=C_LIGHT)
        y += 0.85

    bot_line(slide, "Pre-event phase  |  S3→S8 anticipatory window  |  S8 raised at 14:20 Sep 23  |  Mina (no S8) = no midday dip")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: DURING — The S8 Flip
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide6(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 6, "RQ1 · RQ2", "During Event: The S8 Flip — From Congestion to Clearance", C_YELLOW)

    # Main spatial figure (left)
    rect(slide, 0.2, 0.85, 7.2, 4.6, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图09_空间偏差地理分布图.png", 0.25, 0.90, 7.1, 4.50)
    txt(slide, "A  Spatial distribution of speed deviation during typhoon (S3+, MIDDAY)",
        0.35, 0.87, 6.8, 0.25, size=8, color=C_ACCENT, bold=True)

    # Right: The S8 flip story
    rect(slide, 7.65, 0.85, 5.45, 4.6, fill=C_MID, line=C_ACCENT)
    txt(slide, "The S8 Flip", 7.85, 0.95, 5.0, 0.32, size=13, color=C_YELLOW, bold=True)

    # Timeline boxes
    timeline = [
        ("13:00 Pre-S8", "Mean dev = −0.010\n28.3% roads slower\nSystem-wide congestion", C_RED),
        ("14:20 S8 ↑", "Signal 8 raised\nBehavioural trigger", C_YELLOW),
        ("19:30 Post-S8", "Mean dev = +0.045\n<8% roads slower\nNetwork clears rapidly", C_ACCENT),
    ]
    y = 1.40
    for title, desc, c in timeline:
        rect(slide, 7.85, y, 5.0, 0.82, fill=C_DARK, line=c)
        txt(slide, title, 7.95, y+0.03, 4.8, 0.24, size=10, color=c, bold=True)
        txt(slide, desc, 7.95, y+0.26, 4.8, 0.52, size=8.5, color=C_LIGHT)
        y += 0.92

    # Key number
    rect(slide, 7.85, 4.22, 5.0, 0.55, fill=C_DARK, line=C_YELLOW)
    txt(slide, "Δ = +0.055 in 5 hours", 7.95, 4.27, 4.8, 0.22,
        size=14, color=C_YELLOW, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, "largest 5-hour swing in the entire dataset", 7.95, 4.50, 4.8, 0.22,
        size=8, color=C_GRAY, align=PP_ALIGN.CENTER)

    # Bottom conclusions
    rect(slide, 0.2, 5.60, 12.9, 1.55, fill=C_MID, line=C_ACCENT)
    txt(slide, "Conclusions", 0.45, 5.68, 3.0, 0.30, size=12, color=C_YELLOW, bold=True)

    concl = [
        "S8 is the single most powerful behavioural intervention — the network flips from congestion to clearance within hours",
        "Spatial pattern confirms demand-side dominance: urban core (jobs) clears most, residential periphery mixed",
        "Motorways and trunks are most responsive; local streets retain residual activity throughout the event",
    ]
    y = 6.05
    for c_text in concl:
        rect(slide, 0.45, y+0.03, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, c_text, 0.65, y, 12.0, 0.30, size=9, color=C_LIGHT)
        y += 0.35

    bot_line(slide, "During-event phase  |  S8 triggers +0.055 swing  |  Network clears in ~2h  |  Urban core leads clearance")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: DURING — Speed Polarization
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide7(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 7, "RQ1", "During Event: Speed Polarization — Not All Roads Clear", C_YELLOW)

    # Left: Road category boxplot
    rect(slide, 0.2, 0.85, 6.3, 3.7, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图05_道路类别偏差箱线图.png", 0.25, 0.90, 6.2, 3.60)
    txt(slide, "A  Speed deviation by road category: motorways shift most, local streets mixed",
        0.35, 0.87, 6.0, 0.25, size=8, color=C_ACCENT, bold=True)

    # Right: Signal gradient
    rect(slide, 6.85, 0.85, 6.25, 3.7, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图12_各信号等级五分位偏差.png", 6.90, 0.90, 6.15, 3.60)
    txt(slide, "B  Deviation quantiles by signal level: spread widens with severity",
        7.0, 0.87, 6.0, 0.25, size=8, color=C_ACCENT, bold=True)

    # Bottom: Key findings
    rect(slide, 0.2, 4.70, 12.9, 2.35, fill=C_MID, line=C_ACCENT)
    txt(slide, "Key Findings: The Network Splits into Two Groups", 0.45, 4.78, 12.0, 0.30,
        size=12, color=C_YELLOW, bold=True)

    # Two columns of findings
    col1_x = 0.45
    col2_x = 6.85

    left_findings = [
        ("Demand collapse on arterials",
         "33% of roads clearly faster (>+0.05) at S10\nMotorways/trunks: strongest positive shift\nFewer vehicles → higher observed speeds"),
        ("Residual activity on local streets",
         "11% clearly slower (<-0.05) at S10\nTertiary/residential streets hold activity\nLocal trips persist despite typhoon"),
    ]
    y = 5.20
    for title, desc in left_findings:
        rect(slide, col1_x, y+0.03, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, title, col1_x+0.18, y, 6.0, 0.22, size=10, color=C_ACCENT, bold=True)
        txt(slide, desc, col1_x+0.18, y+0.22, 6.0, 0.45, size=8.5, color=C_LIGHT)
        y += 0.72

    right_findings = [
        ("Signal gradient",
         "IQR: 0.02 (S1) → 0.09 (S10)\n4.5× spread increase from S1 to S10\nEach upgrade widens the distribution"),
        ("Bimodal, not uniform",
         "Roads don't shift uniformly faster\nTwo clusters emerge: faster arterials\nand slower local streets"),
    ]
    y = 5.20
    for title, desc in right_findings:
        rect(slide, col2_x, y+0.03, 0.06, 0.06, fill=C_YELLOW)
        txt(slide, title, col2_x+0.18, y, 5.8, 0.22, size=10, color=C_YELLOW, bold=True)
        txt(slide, desc, col2_x+0.18, y+0.22, 5.8, 0.45, size=8.5, color=C_LIGHT)
        y += 0.72

    bot_line(slide, "Polarization  |  Higher signal → wider spread  |  Demand collapse ≠ uniform  |  Bimodal distribution emerges")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: DURING — Road Disappearance
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide8(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 8, "RQ1 · RQ2", "During Event: Road Disappearance — Supply or Demand?", C_YELLOW)

    # Main figure
    rect(slide, 0.2, 0.85, 8.5, 4.2, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图20_路段韧性地图.png", 0.25, 0.90, 8.4, 4.10)

    # Right: Stats + Interpretation
    rect(slide, 8.95, 0.85, 4.15, 4.2, fill=C_MID, line=C_ACCENT)
    txt(slide, "Coverage Loss by Phase", 9.10, 0.95, 3.8, 0.32, size=12, color=C_YELLOW, bold=True)

    stats = [
        ("Pre-event", "1–3%", "Normal data fluctuation"),
        ("S1–S3", "5–8%", "Minor reduction"),
        ("S8+", "12–18%", "Significant loss"),
        ("S10 peak", "~20%", "Extreme: 1 in 5 roads gone"),
    ]
    y = 1.40
    for phase, pct, desc in stats:
        rect(slide, 9.10, y+0.04, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, phase, 9.28, y, 1.2, 0.22, size=9, color=C_ACCENT, bold=True)
        txt(slide, pct, 10.30, y, 0.8, 0.22, size=11, color=C_YELLOW, bold=True)
        txt(slide, desc, 9.28, y+0.20, 3.5, 0.18, size=7.5, color=C_LIGHT)
        y += 0.62

    # Interpretation
    txt(slide, "Interpretation", 9.10, 4.15, 3.8, 0.28, size=11, color=C_YELLOW, bold=True)
    txt(slide, "Disappearance is primarily demand-side:\n"
         "• Zero cars → TomTom produces no data\n"
         "• Consistent with speed increase on\n"
         "  remaining roads (survivorship)\n"
         "• Local/residential streets most affected\n"
         "• Some supply-side closures at S10 peak",
        9.10, 4.42, 3.8, 0.85, size=8, color=C_LIGHT)

    # Bottom: Coverage comparison
    rect(slide, 0.2, 5.20, 12.9, 1.85, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图16_路网覆盖对比图.png", 0.35, 5.28, 6.15, 1.70)
    txt(slide, "Coverage: Baseline vs Typhoon", 0.40, 5.22, 3.0, 0.18, size=7.5, color=C_GRAY, italic=True)
    img(slide, f"{THESIS}/图17_数据覆盖完整度图.png", 6.80, 5.28, 6.15, 1.70)
    txt(slide, "Coverage completeness over time", 6.85, 5.22, 3.0, 0.18, size=7.5, color=C_GRAY, italic=True)

    bot_line(slide, "Road disappearance  |  Peak ~20% at S10  |  Primarily demand-driven (no cars)  |  Local streets most affected")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: POST-EVENT — Recovery
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide9(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 9, "RQ4", "Post-Event: Rapid Recovery — No Rebound", C_YELLOW)

    # Main: 3-typhoon recovery overview
    rect(slide, 0.2, 0.85, 8.5, 4.0, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图06_三台风恢复动态全图.png", 0.25, 0.90, 8.4, 3.90)

    # Right: Key metrics
    rect(slide, 8.95, 0.85, 4.15, 4.0, fill=C_MID, line=C_ACCENT)
    txt(slide, "Recovery Metrics", 9.10, 0.95, 3.8, 0.32, size=12, color=C_YELLOW, bold=True)

    metrics = [
        ("~100 min", "Return to within ±0.003\nof baseline after All Clear"),
        ("No rebound", "No post-typhoon congestion\nsurge — unlike hurricane lit."),
        ("All 3 typhoons", "Mina, Yagiasha, Madum all\nshow consistent rapid recovery"),
        ("Demand-driven", "Recovery is traffic returning,\nnot infrastructure repair"),
    ]
    y = 1.40
    for num, desc in metrics:
        rect(slide, 9.10, y+0.03, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, num, 9.28, y, 3.5, 0.22, size=10, color=C_ACCENT, bold=True)
        txt(slide, desc, 9.28, y+0.20, 3.5, 0.48, size=8, color=C_LIGHT)
        y += 0.72

    txt(slide, "→ Recovery speed confirms demand-side\n  dominance: if infrastructure damage\n  were primary, recovery would take days",
        9.10, 4.25, 3.8, 0.70, size=8, color=C_GRAY, italic=True)

    # Bottom detail
    rect(slide, 0.2, 5.0, 12.9, 2.05, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图07_叶加沙恢复期细节.png", 0.35, 5.08, 6.15, 1.85)
    txt(slide, "Yagiasha recovery detail", 0.40, 5.02, 3.0, 0.16, size=7.5, color=C_GRAY, italic=True)
    img(slide, f"{THESIS}/图08_道路类别恢复曲线.png", 6.80, 5.08, 6.15, 1.85)
    txt(slide, "Recovery by road category: motorways fastest", 6.85, 5.02, 3.0, 0.16, size=7.5, color=C_GRAY, italic=True)

    bot_line(slide, "Post-event phase  |  ~100 min recovery  |  No rebound  |  Consistent across all 3 typhoons  |  Demand-driven")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: GRID REGRESSION
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide10(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 10, "RQ3", "Grid Regression: What Determines Speed Change?", C_ACCENT)

    # Coefficient chart
    rect(slide, 0.2, 0.85, 7.2, 5.2, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图46_回归系数图.png", 0.25, 0.90, 7.1, 5.10)

    # Right: Model summary
    rect(slide, 7.65, 0.85, 5.45, 5.2, fill=C_MID, line=C_ACCENT)

    txt(slide, "Model Summary", 7.85, 0.95, 5.0, 0.35, size=13, color=C_YELLOW, bold=True)

    model_info = [
        ("Y variable", "Length-weighted % of roads with\ndeviation > +0.05 (\"clearly better\")"),
        ("Unit", "500m × 500m grid cell\n(≥3 roads, ≥500m total length)"),
        ("X variables", "10 POI log-densities + road structure\n+ demographics + signal dummies"),
        ("Best model", "PM_PEAK: adj-R² = 0.188\nMIDDAY: adj-R² = 0.136"),
    ]
    y = 1.40
    for title, desc in model_info:
        txt(slide, title, 7.85, y, 5.0, 0.22, size=9, color=C_ACCENT, bold=True)
        txt(slide, desc, 7.85, y+0.20, 5.0, 0.40, size=8.5, color=C_LIGHT)
        y += 0.68

    # Key predictors
    txt(slide, "Consistent Predictors", 7.85, 4.10, 5.0, 0.30, size=11, color=C_YELLOW, bold=True)

    predictors = [
        ("+ Intersection degree", "More connected roads → more improvement"),
        ("+ Transport POI", "Near transit hubs → more roads get faster"),
        ("− Road density", "Denser network → less improvement"),
        ("− Signal 8", "S8 suppresses improvement (pre-S8 rush)"),
    ]
    y = 4.45
    for pred, desc in predictors:
        rect(slide, 7.85, y+0.04, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, pred, 8.00, y, 4.8, 0.22, size=9, color=C_ACCENT, bold=True)
        txt(slide, desc, 8.00, y+0.18, 4.8, 0.20, size=8, color=C_LIGHT)
        y += 0.42

    # Bottom
    rect(slide, 0.2, 6.20, 12.9, 0.85, fill=C_MID, line=C_ACCENT)
    txt(slide, "Interpretation:  Grids with complex road networks near transit hubs clear fastest during typhoons.  "
         "Dense road networks in commercial areas show less improvement — residual local activity.  "
         "Demographic variables have limited explanatory power at grid scale.",
        0.45, 6.35, 12.3, 0.55, size=9.5, color=C_LIGHT)

    bot_line(slide, "Regression analysis  |  500m grid  |  Ragasa S3+  |  Length-weighted aggregation")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: CASE STUDY — Grid Visualization
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide11(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 11, "RQ3", "Case Study: Single Grid × Regression Variables", C_ACCENT)

    # Left: Grid visualization
    rect(slide, 0.2, 0.85, 6.5, 3.8, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图41_网格可视化.png", 0.25, 0.90, 6.4, 3.70)
    txt(slide, "A  Mong Kok grid: 9 roads, pct_better=0.83, mean_dev=+0.154",
        0.35, 0.87, 6.0, 0.25, size=8, color=C_ACCENT, bold=True)

    # Right: Regression schematic
    rect(slide, 6.95, 0.85, 6.15, 3.8, fill=C_MID, line=C_ACCENT)
    img(slide, f"{THESIS}/图37_回归示意图.png", 7.00, 0.90, 6.05, 3.70)
    txt(slide, "B  Wan Chai study road: 861 POIs in 500m buffer",
        7.10, 0.87, 6.0, 0.25, size=8, color=C_ACCENT, bold=True)

    # Bottom: Grid story
    rect(slide, 0.2, 4.80, 12.9, 2.25, fill=C_MID, line=C_ACCENT)
    txt(slide, "From Road to Grid: Why Aggregation Works", 0.45, 4.90, 12.0, 0.30,
        size=13, color=C_YELLOW, bold=True)

    points = [
        "Road-level R² ≈ 0.02–0.06 → Grid-level R² ≈ 0.14–0.19: spatial aggregation reduces micro-noise",
        "500m grid captures neighborhood-scale determinants (POI, demographics) better than single roads",
        "Length-weighting ensures long arterials contribute proportionally more to grid metrics",
        "The Mong Kok case: high intersection density + transport POI + moderate road density → 83% roads clearly faster",
    ]
    y = 5.30
    for p in points:
        rect(slide, 0.45, y+0.03, 0.06, 0.06, fill=C_ACCENT)
        txt(slide, p, 0.65, y, 12.0, 0.30, size=9, color=C_LIGHT)
        y += 0.37

    bot_line(slide, "Grid aggregation  |  500m cells  |  Length-weighted  |  Neighborhood-scale determinants")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: CONCLUSIONS
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide12(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    top_bar(slide, 12, "SUMMARY", "Conclusions & Contributions", C_ACCENT)

    # Three columns
    phases = [
        ("Pre-Event", "#4fc3f7",
         ["S8 is the behavioural trigger (not S1/S3)",
          "Mina (no S8): midday speeds elevated; Yagiasha: dip",
          "Pre-S8 congestion surge: −0.010 at 13:00",
          "System-wide last-minute rush (commute, not shopping)"]),
        ("During Event", "#ef5350",
         ["S8 triggers rapid network clearance (+0.045 in 5h)",
          "Peak deviation +0.046 at S10 morning",
          "Speed polarization: 33% clearly faster, 11% slower",
          "~20% roads disappear at S10 peak",
          "Motorways/trunks most affected; local streets mixed"]),
        ("Post-Event", "#81c784",
         ["Recovery within ~100 min of All Clear",
          "No post-typhoon rebound congestion",
          "Recovery is demand-driven (not supply-constrained)",
          "Consistent pattern across all 3 typhoons"]),
    ]

    for i, (phase, color, points) in enumerate(phases):
        x = 0.25 + i * 4.35
        rect(slide, x, 0.95, 4.10, 4.5, fill=C_MID, line=color)

        # Phase header
        rect(slide, x, 0.95, 4.10, 0.45, fill=color)
        txt(slide, phase, x+0.15, 0.98, 3.8, 0.38, size=16, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)

        y = 1.55
        for p in points:
            rect(slide, x+0.15, y+0.03, 0.05, 0.05, fill=color)
            txt(slide, p, x+0.30, y, 3.6, 0.55, size=9, color=C_LIGHT)
            y += 0.62

    # Contributions
    rect(slide, 0.2, 5.60, 5.8, 1.45, fill=C_MID, line=C_ACCENT)
    txt(slide, "Contributions", 0.40, 5.68, 5.4, 0.30, size=12, color=C_YELLOW, bold=True)
    contribs = [
        "First empirical evidence that Signal 8 (not S1/S3) is the behavioural tipping point for pre-typhoon mobility",
        "500m grid regression framework linking road performance to urban structure",
        "Three-phase (pre/during/post) resilience analysis with per-slot granularity",
    ]
    y = 6.05
    for c in contribs:
        rect(slide, 0.40, y+0.02, 0.05, 0.05, fill=C_ACCENT)
        txt(slide, c, 0.55, y, 5.2, 0.30, size=8.5, color=C_LIGHT)
        y += 0.32

    # Limitations
    rect(slide, 6.25, 5.60, 6.85, 1.45, fill=C_MID, line=C_ACCENT)
    txt(slide, "Limitations & Future Work", 6.45, 5.68, 6.4, 0.30, size=12, color=C_YELLOW, bold=True)
    lims = [
        "Data gaps (Sep 22) limit pre-S1 analysis for some typhoons",
        "Grid-level adj-R² ~0.19: substantial unexplained variance remains",
        "TomTom data excludes zero-traffic roads → survivorship bias",
        "Future: traffic flow volume data, incident integration, causal identification",
    ]
    y = 6.05
    for l in lims:
        rect(slide, 6.45, y+0.02, 0.05, 0.05, fill=C_GRAY)
        txt(slide, l, 6.60, y, 6.2, 0.30, size=8.5, color=C_LIGHT)
        y += 0.32

    bot_line(slide, "Conclusions  |  Pre-during-post resilience framework  |  Empirical evidence from Hong Kong")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
def make_slide13(prs):
    slide = new_slide(prs)
    rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)
    rect(slide, 3.0, 3.0, 7.33, 0.06, fill=C_ACCENT)
    txt(slide, "Thank You", 3.0, 2.0, 7.33, 1.0, size=42, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, "Questions & Discussion", 3.0, 3.3, 7.33, 0.7, size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)
    txt(slide, "MPhil Thesis Defense  ·  Department of Urban Planning and Design  ·  HKU  ·  2026",
        3.0, 4.5, 7.33, 0.5, size=12, color=C_GRAY, align=PP_ALIGN.CENTER)
    rect(slide, 0, 7.28, 13.33, 0.04, fill=C_ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# GENERATE ALL
# ═══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    for i, make_fn in enumerate([
        make_slide1, make_slide2, make_slide3, make_slide4, make_slide5,
        make_slide6, make_slide7, make_slide8, make_slide9, make_slide10,
        make_slide11, make_slide12, make_slide13
    ], start=1):
        print(f"Slide {i:02d}...", end=" ", flush=True)
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        # Remove default slide
        for s in list(prs.slides):
            prs.slides._sldIdLst.remove(s._element)
        make_fn(prs)
        out = f"{THESIS}/答辩PPT_slide{i:02d}.pptx"
        prs.save(out)
        print(f"Saved: {out}")

    print("\nAll 13 slides generated!")

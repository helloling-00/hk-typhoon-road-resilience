"""
答辩PPT Slide 5 — Pre-Event: S8 Is the True Behavioural Trigger
Mina (no S8) vs Yagiasha (S8+) speed shape comparison
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

THESIS = "/Users/helloling/workspace/thesis"

C_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
C_MID    = RGBColor(0x16, 0x21, 0x3e)
C_ACCENT = RGBColor(0x0f, 0x8b, 0x8d)
C_LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
C_YELLOW = RGBColor(0xff, 0xc8, 0x00)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_GRAY   = RGBColor(0xaa, 0xaa, 0xaa)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill: shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else: shape.fill.background()
    if line: shape.line.color.rgb = line; shape.line.width = Pt(1)
    else: shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, color=C_WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    return txb

def img(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

# ── Background ────────────────────────────────────────────────────────────────
rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)

# ── Top bar ────────────────────────────────────────────────────────────────────
rect(slide, 0, 0, 13.33, 0.75, fill=C_MID)
txt(slide, "05", 0.15, 0.12, 0.5, 0.5, size=11, color=C_GRAY)
rect(slide, 0.5, 0.17, 2.0, 0.40, fill=C_YELLOW)
txt(slide, "RQ1 · RQ2", 0.55, 0.18, 1.9, 0.38,
    size=10, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)
txt(slide, "Pre-Event: S8 Is the True Behavioural Trigger",
    2.7, 0.10, 8.5, 0.58, size=22, color=C_WHITE, bold=True)
rect(slide, 11.4, 0.15, 1.7, 0.42, fill=C_YELLOW)
txt(slide, "PRE-TYPHOON", 11.42, 0.17, 1.66, 0.38,
    size=10, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)

# ── MAIN FIGURE: Mina vs Yagiasha (图47) — full width ──────────────────────────
FIG_TOP = 0.85
rect(slide, 0.2, FIG_TOP, 12.9, 3.15, fill=C_MID, line=C_ACCENT)
img(slide, f"{THESIS}/图47_速度形状对比_米娜vs叶加沙.png", 0.25, FIG_TOP+0.05, 12.8, 3.05)

txt(slide, "A  Mina (no S8): all-day elevated speeds  —  Yagiasha (S8+): midday dip ONLY on Sep 23 (pre-S8)",
    0.35, FIG_TOP+0.02, 12.0, 0.25, size=9, color=C_ACCENT, bold=True)

# ── BOTTOM ROW: Two panels ────────────────────────────────────────────────────
BOT_TOP = FIG_TOP + 3.30

# Left: Detailed pre-S8 surge (图43b)
rect(slide, 0.2, BOT_TOP, 6.3, 3.0, fill=C_MID, line=C_ACCENT)
img(slide, f"{THESIS}/图43b_preS8_surge.png", 0.25, BOT_TOP+0.05, 6.2, 2.90)
txt(slide, "B  Yagiasha Sep 23: Pre-S8 dip → S8 clearance in detail",
    0.35, BOT_TOP+0.02, 6.0, 0.25, size=9, color=C_ACCENT, bold=True)

# Right: Key evidence panel
rect(slide, 6.85, BOT_TOP, 6.25, 3.0, fill=C_MID, line=C_ACCENT)
txt(slide, "Key Evidence", 7.05, BOT_TOP+0.08, 5.8, 0.30, size=13, color=C_YELLOW, bold=True)

evidence = [
    ("Mina (max S3, no S8)",
     "Sep 18–19 midday: speeds +0.005 to +0.015 above baseline\n"
     "→ No congestion surge; demand already suppressed"),
    ("Yagiasha (S8 raised 14:20)",
     "Sep 23 midday: speed dips to −0.010 at 13:00\n"
     "→ 31.3% roads slower, system-wide last-minute rush"),
    ("5-hour swing",
     "−0.010 (13:00) → +0.045 (19:30) after S8\n"
     "→ S8 is the behavioural tipping point, not S1 or S3"),
]
y = BOT_TOP + 0.50
for title, desc in evidence:
    rect(slide, 7.05, y+0.03, 0.07, 0.07, fill=C_ACCENT)
    txt(slide, title, 7.25, y, 5.5, 0.25, size=10, color=C_ACCENT, bold=True)
    txt(slide, desc, 7.25, y+0.22, 5.5, 0.55, size=9, color=C_LIGHT)
    y += 0.85

# ── Bottom line ────────────────────────────────────────────────────────────────
rect(slide, 0, 7.28, 13.33, 0.04, fill=C_ACCENT)
txt(slide, "Pre-event phase  |  S3→S8 anticipatory window  |  S8 raised at 14:20  |  Mina (no S8) = no midday dip",
    0.3, 7.33, 12.0, 0.20, size=8, color=C_GRAY, italic=True)

out = f"{THESIS}/答辩PPT_slide05.pptx"
prs.save(out)
print(f"Saved: {out}")

"""
答辩PPT Slide 6 — Pre-Event Behavioural Anomalies
Two key findings: (1) Anticipatory mobility before S1, (2) Pre-S8 congestion surge
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

THESIS = "/Users/helloling/workspace/thesis"

C_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
C_MID    = RGBColor(0x16, 0x21, 0x3e)
C_ACCENT = RGBColor(0x0f, 0x8b, 0x8d)
C_LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
C_YELLOW = RGBColor(0xff, 0xc8, 0x00)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_GRAY   = RGBColor(0xaa, 0xaa, 0xaa)
C_RED    = RGBColor(0xe0, 0x52, 0x52)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill: shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else: shape.fill.background()
    if line: shape.line.color.rgb = line; shape.line.width = Pt(1)
    else: shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, color=C_WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    return txb

def img(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

# ── Background ────────────────────────────────────────────────────────────────
rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)

# ── Top bar ────────────────────────────────────────────────────────────────────
rect(slide, 0, 0, 13.33, 0.75, fill=C_MID)
txt(slide, "06", 0.15, 0.12, 0.5, 0.5, size=11, color=C_GRAY)
rect(slide, 0.5, 0.17, 2.0, 0.40, fill=C_ACCENT)
txt(slide, "RQ1 · RQ2", 0.55, 0.18, 1.9, 0.38,
    size=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(slide, "Pre-Event Behavioural Anomalies",
    2.7, 0.10, 8.5, 0.58, size=24, color=C_WHITE, bold=True)
rect(slide, 11.4, 0.15, 1.7, 0.42, fill=C_YELLOW)
txt(slide, "PRE-TYPHOON", 11.42, 0.17, 1.66, 0.38,
    size=10, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)

# ── Layout: top 60% = two figures side by side, bottom 40% = explanations ────
FIG_TOP = 0.85
FIG_H   = 3.55

# ── LEFT: 图01 — Anticipatory mobility ─────────────────────────────────────────
rect(slide, 0.2, FIG_TOP, 6.3, FIG_H, fill=C_MID, line=C_ACCENT)
img(slide, f"{THESIS}/图43a_preS1_anticipatory.png", 0.25, FIG_TOP+0.05, 6.2, FIG_H-0.10)

# Label on figure
txt(slide, "A  Pre-Signal 1: Anticipatory mobility",
    0.3, FIG_TOP+0.04, 5.0, 0.28, size=9, color=C_ACCENT, bold=True)

# ── RIGHT: 图11 — Pre-S8 congestion surge ──────────────────────────────────────
rect(slide, 6.85, FIG_TOP, 6.25, FIG_H, fill=C_MID, line=C_ACCENT)
img(slide, f"{THESIS}/图43b_preS8_surge.png",
    6.90, FIG_TOP+0.05, 6.15, FIG_H-0.10)

txt(slide, "B  Pre-Signal 8: Last-minute congestion surge",
    6.95, FIG_TOP+0.04, 5.0, 0.28, size=9, color=C_ACCENT, bold=True)

# ── BOTTOM: Key Findings ───────────────────────────────────────────────────────
FIND_TOP = FIG_TOP + FIG_H + 0.12
rect(slide, 0.2, FIND_TOP, 12.9, 2.55, fill=C_MID, line=C_ACCENT)

# Two columns of findings
# Left column: Finding 1
col1_x = 0.45
txt(slide, "Finding 1  ·  Anticipatory demand suppression",
    col1_x, FIND_TOP+0.10, 6.0, 0.32, size=13, color=C_YELLOW, bold=True)

# Key numbers in boxes
metrics_y = FIND_TOP + 0.50
for i, (label, value) in enumerate([
    ("Peak network deviation before S1 (Mina, Sep 17)", "+0.012 at 17:00"),
    ("Roads affected", "Tertiary & secondary roads most responsive"),
    ("Behavioural interpretation", "Early departure, last-minute provisioning"),
]):
    rect(slide, col1_x, metrics_y + i*0.52, 5.8, 0.46, fill=C_DARK, line=C_ACCENT)
    txt(slide, label,  col1_x+0.10, metrics_y+i*0.52+0.03, 5.6, 0.20,
        size=8, color=C_GRAY, italic=False)
    txt(slide, value, col1_x+0.10, metrics_y+i*0.52+0.22, 5.6, 0.22,
        size=10, color=C_WHITE, bold=True)

# Right column: Finding 2
col2_x = 7.10
txt(slide, "Finding 2  ·  Pre-Signal 8 congestion surge",
    col2_x, FIND_TOP+0.10, 6.0, 0.32, size=13, color=C_YELLOW, bold=True)

for i, (label, value) in enumerate([
    ("Pre-S8 dip (Yagiasha, Sep 23)", "Mean deviation drops to −0.010 at 13:00"),
    ("Roads affected", "31.3% of roads slower than workday baseline"),
    ("Behavioural interpretation", "Rush to supermarkets, collect children, commute home"),
]):
    rect(slide, col2_x, metrics_y + i*0.52, 5.8, 0.46, fill=C_DARK, line=C_ACCENT)
    txt(slide, label,  col2_x+0.10, metrics_y+i*0.52+0.03, 5.6, 0.20,
        size=8, color=C_GRAY, italic=False)
    txt(slide, value, col2_x+0.10, metrics_y+i*0.52+0.22, 5.6, 0.22,
        size=10, color=C_WHITE, bold=True)

# ── Bottom line ────────────────────────────────────────────────────────────────
rect(slide, 0, 7.28, 13.33, 0.04, fill=C_ACCENT)
txt(slide, "Pre-event phase  |  Demand begins suppressing before formal warning issuance  |  1.5 min",
    0.3, 7.33, 12.0, 0.20, size=8, color=C_GRAY, italic=True)

out = f"{THESIS}/答辩PPT_slide06.pptx"
prs.save(out)
print(f"Saved: {out}")

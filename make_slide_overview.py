"""
生成答辩PPT：Overview slide — Network-Level Speed Deviation (overview)
概览页：台风全阶段网络速度偏差时序图
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

THESIS = "/Users/helloling/workspace/thesis"

C_DARK   = RGBColor(0x1a, 0x1a, 0x2e)
C_MID    = RGBColor(0x16, 0x21, 0x3e)
C_ACCENT = RGBColor(0x0f, 0x8b, 0x8d)
C_LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
C_YELLOW = RGBColor(0xff, 0xc8, 0x00)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_GRAY   = RGBColor(0xaa, 0xaa, 0xaa)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, color=C_WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.color.rgb = color
    run.font.bold    = bold
    run.font.italic  = italic
    return txb

def img(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

# ── Background ────────────────────────────────────────────────────────────────
rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK)

# ── Top title bar ─────────────────────────────────────────────────────────────
rect(slide, 0, 0, 13.33, 0.75, fill=C_MID)

txt(slide, "04", 0.15, 0.12, 0.5, 0.5, size=11, color=C_GRAY)
rect(slide, 0.5, 0.17, 2.8, 0.40, fill=C_ACCENT)
txt(slide, "RQ1 · RQ2 · RQ3", 0.55, 0.18, 2.7, 0.38,
    size=11, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(slide, "Overview: Network Response across All Typhoon Phases",
    3.5, 0.10, 9.0, 0.58, size=22, color=C_WHITE, bold=True)

# Right tag
rect(slide, 11.4, 0.15, 1.7, 0.42, fill=C_ACCENT)
txt(slide, "ALL PHASES", 11.42, 0.17, 1.66, 0.38,
    size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── Main figure ───────────────────────────────────────────────────────────────
rect(slide, 0.2, 0.85, 12.9, 5.35, fill=C_MID, line=C_ACCENT)
img(slide, f"{THESIS}/图42_台风时序网络偏差.png",
    0.25, 0.90, 12.80, 5.25)

# ── Key takeaway strip ────────────────────────────────────────────────────────
rect(slide, 0.2, 6.28, 12.9, 0.92, fill=C_MID, line=C_ACCENT)

takeaways = [
    ("Pre-event",   "Anticipatory demand suppression begins hours before Signal 1"),
    ("During",      "Network clears rapidly after Signal 8 — +4–5% above workday baseline"),
    ("Post-event",  "Recovery within ~2 hrs of All Clear; no rebound congestion"),
]

x = 0.35
for phase, desc in takeaways:
    rect(slide, x, 6.37, 0.06, 0.06, fill=C_ACCENT)
    txt(slide, phase + ":", x+0.12, 6.33, 1.3, 0.32,
        size=10, color=C_YELLOW, bold=True)
    txt(slide, desc, x+0.12, 6.62, 3.8, 0.28,
        size=9, color=C_LIGHT, italic=False)
    x += 4.3

# ── Bottom line ───────────────────────────────────────────────────────────────
rect(slide, 0, 7.28, 13.33, 0.04, fill=C_ACCENT)
txt(slide, "Ragasa (Sep 2025)  |  30-min TomTom floating-car data  |  ~6,000–7,000 roads per slot",
    0.3, 7.33, 12.0, 0.20, size=8, color=C_GRAY, italic=True)

out = f"{THESIS}/答辩PPT_slide_overview.pptx"
prs.save(out)
print(f"Saved: {out}")
